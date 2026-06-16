#!/usr/bin/env python3
"""Build East Meets Feed – professional PPTX from HTML content"""

import io, re, base64
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ─────────────────────────────────────────────────────────────────
GOLD   = RGBColor(0xC9, 0xA9, 0x6E)
GOLD2  = RGBColor(0xE8, 0xC9, 0x8A)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
DARK2  = RGBColor(0x11, 0x11, 0x11)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xF5, 0xF4, 0xF0)
LGRAY  = RGBColor(0xF7, 0xF7, 0xF5)
MGRAY  = RGBColor(0x99, 0x99, 0x99)
DGRAY  = RGBColor(0x66, 0x66, 0x66)
XGRAY  = RGBColor(0xE8, 0xE8, 0xE8)
BGRAY  = RGBColor(0xFA, 0xFA, 0xFA)
GREEN  = RGBColor(0x2D, 0x7A, 0x2D)
LGREEN = RGBColor(0xE8, 0xF4, 0xE8)
AMBER  = RGBColor(0xB3, 0x62, 0x00)
LAMBER = RGBColor(0xFF, 0xF0, 0xE0)
IGRED  = RGBColor(0xDC, 0x27, 0x43)
XHSRED = RGBColor(0xFE, 0x2C, 0x55)
INDIG  = RGBColor(0x63, 0x66, 0xF1)

# ── Slide geometry ───────────────────────────────────────────────────────────
SW = Inches(13.33)   # slide width
SH = Inches(7.5)     # slide height

i = Inches          # shorthand
p = Pt

BODY_FONT = "Calibri"
HEAD_FONT = "Calibri Light"

# ── Low-level helpers ────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def txt(slide, x, y, w, h, text, size=12, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT, font=BODY_FONT, wrap=True, spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    if spacing:
        para.space_after = Pt(spacing)
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = p(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def txt_multi(slide, x, y, w, h, lines, font=BODY_FONT):
    """lines = list of (text, size, bold, italic, color, align)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, italic, color, align) in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.name = font
        run.font.size = p(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def label(slide, x, y, w, h, text, color=MGRAY):
    return txt(slide, x, y, w, h, text.upper(), size=8, bold=True, color=color)


def pic(slide, img_buf, x, y, w, h=None):
    img_buf.seek(0)
    if h:
        return slide.shapes.add_picture(img_buf, x, y, w, h)
    else:
        return slide.shapes.add_picture(img_buf, x, y, w)


def slide_num_tag(slide, n, total=10):
    txt(slide, SW - i(1.6), SH - i(0.45), i(1.4), i(0.3),
        f"{n:02d} / {total:02d}", size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


def project_tag(slide, dark=False):
    c = RGBColor(0xBB, 0xBB, 0xBB) if dark else MGRAY
    txt(slide, SW - i(2.8), i(0.18), i(2.6), i(0.3),
        "EAST MEETS FEED", size=8, bold=True, color=c, align=PP_ALIGN.RIGHT)


def accent_bar(slide):
    rect(slide, 0, 0, SW, i(0.07), GOLD)


def section_header(slide, eyebrow, title, subtitle=None):
    accent_bar(slide)
    y = i(0.35)
    txt(slide, i(0.7), y, i(10), i(0.25),
        eyebrow.upper(), size=9, bold=True, color=GOLD)
    txt(slide, i(0.7), y + i(0.3), i(10), i(0.55),
        title, size=26, bold=True, color=DARK, font=HEAD_FONT)
    if subtitle:
        txt(slide, i(0.7), y + i(0.9), i(10), i(0.3),
            subtitle, size=12, italic=True, color=MGRAY)


# ── Extract images ───────────────────────────────────────────────────────────

def extract_images(html_path):
    with open(html_path, 'r') as f:
        content = f.read()
    b64_list = re.findall(r'src="(data:image/png;base64,[^"]+)"', content)
    bufs = []
    for b64_str in b64_list:
        data = base64.b64decode(b64_str.split(',')[1])
        img = Image.open(io.BytesIO(data))
        if img.mode == 'RGBA':
            bg = Image.new('RGB', img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            img = bg
        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=85)
        buf.seek(0)
        bufs.append(buf)
    return bufs


# ── Slide builders ──────────────────────────────────────────────────────────

def slide1_title(prs, imgs):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    set_bg(s, DARK)

    # Gold top bar
    rect(s, 0, 0, SW, i(0.07), GOLD)

    # Left panel content area
    lw = i(6.4)
    lx = i(0.85)

    txt(s, lx, i(0.9), i(5), i(0.25),
        "A CROSS-PLATFORM BRAND ACTIVATION",
        size=9, bold=True, color=GOLD)

    # Big title
    txt(s, lx, i(1.3), i(5.5), i(3.0),
        "East\nMeets\nFeed",
        size=64, bold=True, color=WHITE, font=HEAD_FONT)

    # Gold divider
    rect(s, lx, i(4.45), i(0.6), i(0.04), GOLD)

    # Subtitle
    txt(s, lx, i(4.6), i(5.5), i(0.55),
        "Building @Yujia's Presence on Instagram · Douyin · Xiaohongshu",
        size=13, color=RGBColor(0xAA, 0xAA, 0xAA))

    # Platform chips row
    chips = [("IG", IGRED), ("Douyin", DARK2), ("XHS", XHSRED)]
    cx = lx
    cy = i(5.3)
    for label_text, chip_color in chips:
        rect(s, cx, cy, i(1.1), i(0.3), chip_color)
        txt(s, cx + i(0.07), cy + i(0.03), i(1.0), i(0.26),
            label_text, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += i(1.22)

    # Meta line
    txt(s, lx, i(5.78), i(6), i(0.3),
        "Eric Wang  ·  Social Media & Brand  ·  NYU  ·  June 2026",
        size=10, color=RGBColor(0x66, 0x66, 0x66))

    # Right photo – img[0] (1320×1761 portrait)
    ph = i(6.8)
    pw = ph * (1320 / 1761)
    px = SW - pw - i(0.2)
    py = i(0.35)
    pic(s, imgs[0], px, py, pw, ph)

    # Gradient fade over left edge of photo (dark rectangle with transparency using XML)
    # Simulate with a narrow dark rect on the seam
    rect(s, px - i(0.5), 0, i(0.9), SH, DARK)

    slide_num_tag(s, 1)
    return s


def slide2_subject(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    lx = i(0.7)
    label(s, lx, i(0.45), i(5), i(0.22), "02  /  The Subject & Why")

    txt(s, lx, i(0.85), i(6.5), i(1.3),
        "She takes photos that\ndeserve a bigger audience.",
        size=24, bold=True, color=DARK, font=HEAD_FONT)
    # gold "bigger audience" highlight effect – overlay with gold text on same area
    txt(s, lx + i(2.22), i(1.37), i(4), i(0.5),
        "bigger audience.", size=24, bold=True, color=GOLD, font=HEAD_FONT)

    y = i(2.4)
    label(s, lx, y, i(4), i(0.22), "Who is Yujia")
    txt(s, lx, y + i(0.28), i(6.2), i(0.8),
        "Yujia Xie — Data Engineer Tech Lead at CVS, NYC.  B.S. Data Science, UCSD  ·  M.S. Data Science, Columbia.",
        size=12, color=DGRAY)

    # Tags
    tags = ["Lifestyle", "Fashion", "NYC", "Cross-Cultural"]
    tx = lx
    ty = y + i(1.2)
    for tag in tags:
        tw = i(len(tag) * 0.09 + 0.5)
        rect(s, tx, ty, tw, i(0.26), CREAM)
        txt(s, tx + i(0.1), ty + i(0.02), tw - i(0.1), i(0.24),
            tag, size=10, bold=True, color=RGBColor(0xA0, 0x78, 0x40))
        tx += tw + i(0.12)

    y2 = i(4.0)
    label(s, lx, y2, i(4), i(0.22), "The Opportunity")
    txt(s, lx, y2 + i(0.3), i(6.2), i(0.55),
        "2,900+ combined followers across three platforms — but no consistent strategy, no cross-posting, no growth system.",
        size=12, color=DGRAY)

    y3 = i(5.2)
    rect(s, lx, y3, i(4.2), i(0.42), DARK)
    txt(s, lx + i(0.2), y3 + i(0.08), i(4.0), i(0.3),
        "STRATEGIST  ·  ACTIVATOR  ·  PRODUCER",
        size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Right photo – img[1] (921×2000 portrait)
    ph = i(6.8)
    pw = ph * (921 / 2000)
    px = SW - pw - i(0.15)
    pic(s, imgs[1], px, i(0.35), pw, ph)

    slide_num_tag(s, 2)
    return s


def slide3_strategy(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, LGRAY)
    accent_bar(s)
    project_tag(s)

    section_header(s, "03  /  Platform Strategy",
                   "One Creator. Three Platforms. One Brand.",
                   "Mapping @Yujia's content to the right audience on each platform")

    # Three platform columns
    cols = [
        {
            "name": "Instagram", "handle": "@yujiaaa_xie",
            "audience": "English-speaking Western market · NYC creatives & lifestyle",
            "content": "Polished lifestyle photography · editorial fashion · cultural spaces",
            "tone": "Aspirational · Quiet luxury · Authentic",
            "chip": "2,387 followers — strongest base", "chip_color": LGREEN, "chip_text": GREEN,
            "icon_color": IGRED,
        },
        {
            "name": "Douyin", "handle": "@Xyjjj",
            "audience": "Chinese diaspora + mainland China · algorithm-driven discovery",
            "content": "Emotional lifestyle · poetic captions · relatable slice-of-life",
            "tone": "Poetic · Resonant · Emotionally direct",
            "chip": "Highest viral upside", "chip_color": LAMBER, "chip_text": AMBER,
            "icon_color": DARK,
        },
        {
            "name": "Xiaohongshu", "handle": "谢雨佳 · 小红书",
            "audience": "Chinese women · lifestyle & fashion discovery · aspiration-driven",
            "content": "Curated photo posts · location tags · OOTD & cultural experiences",
            "tone": "Aspirational · Discovery-driven · Community",
            "chip": "530 followers — high engagement", "chip_color": RGBColor(0xE8, 0xEE, 0xFF), "chip_text": INDIG,
            "icon_color": XHSRED,
        },
    ]

    col_w = i(3.7)
    gap = i(0.4)
    total_w = col_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    card_y = i(1.65)
    card_h = i(5.1)

    for ci, col in enumerate(cols):
        cx = start_x + ci * (col_w + gap)
        rect(s, cx, card_y, col_w, card_h, WHITE)

        # Colored top strip
        rect(s, cx, card_y, col_w, i(0.08), col["icon_color"])

        inner_x = cx + i(0.22)
        y = card_y + i(0.25)

        # Platform name
        txt(s, inner_x, y, col_w - i(0.4), i(0.38),
            col["name"], size=17, bold=True, color=DARK, font=HEAD_FONT)
        txt(s, inner_x, y + i(0.38), col_w - i(0.4), i(0.22),
            col["handle"], size=10, color=MGRAY)

        # Separator
        rect(s, inner_x, y + i(0.68), col_w - i(0.44), i(0.01), XGRAY)

        y2 = y + i(0.85)
        for field, value in [("Audience", col["audience"]),
                              ("Content", col["content"]),
                              ("Tone", col["tone"])]:
            label(s, inner_x, y2, col_w - i(0.44), i(0.2), field)
            txt(s, inner_x, y2 + i(0.22), col_w - i(0.44), i(0.6),
                value, size=11, color=DGRAY)
            y2 += i(1.0)

        # Chip
        cw = col_w - i(0.44)
        rect(s, inner_x, card_y + card_h - i(0.55), cw, i(0.3), col["chip_color"])
        txt(s, inner_x + i(0.1), card_y + card_h - i(0.52), cw - i(0.1), i(0.27),
            col["chip"], size=10, bold=True, color=col["chip_text"])

    # Target bar
    bar_y = SH - i(0.52)
    rect(s, i(0.5), bar_y - i(0.15), SW - i(1.0), i(0.42), DARK)
    txt(s, i(0.8), bar_y - i(0.1), i(4), i(0.3),
        "TARGET AUDIENCE", size=8, bold=True, color=RGBColor(0x55, 0x55, 0x55))
    txt(s, i(3.2), bar_y - i(0.1), i(5), i(0.3),
        "Women  ·  Ages 24–35  ·  Lifestyle, Fashion & Culture",
        size=11, bold=True, color=GOLD)

    slide_num_tag(s, 3)
    return s


def slide4_plan(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    section_header(s, "04  /  The Activation", "Content-First. Purely Organic.")
    accent_bar(s)
    project_tag(s)

    # Left strategy column
    lx = i(0.7)
    lw = i(4.0)
    y = i(1.7)
    label(s, lx, y - i(0.25), lw, i(0.22), "Strategy Pillars")

    pillars = [
        ("Pillar 1 — Aesthetic Visuals",
         "High-quality photography in iconic locations. Fashion-forward, editorial feel. Consistent visual identity across all platforms."),
        ("Pillar 2 — Resonant Captions",
         "Poetic, introspective copy that invites emotional response — or open-ended questions that pull audiences into conversation."),
    ]
    for title_p, desc_p in pillars:
        rect(s, lx, y, i(0.05), i(0.75), GOLD)
        rect(s, lx + i(0.05), y, lw - i(0.05), i(0.75), CREAM)
        txt(s, lx + i(0.18), y + i(0.06), lw - i(0.3), i(0.25),
            title_p, size=12, bold=True, color=DARK)
        txt(s, lx + i(0.18), y + i(0.33), lw - i(0.3), i(0.45),
            desc_p, size=10, color=DGRAY)
        y += i(0.92)

    # No hashtag box
    y += i(0.1)
    rect(s, lx, y, lw, i(0.7), RGBColor(0xFF, 0xF8, 0xF0))
    # Dashed border effect (solid thin line)
    border = slide.shapes.add_shape(1, lx, y, lw, i(0.7)) if False else None
    txt(s, lx + i(0.15), y + i(0.06), lw - i(0.2), i(0.22),
        "🚫  No Hashtags — Intentional", size=11, bold=True, color=AMBER)
    txt(s, lx + i(0.15), y + i(0.3), lw - i(0.2), i(0.38),
        "Past experience showed hashtags yielded no meaningful reach. Phase 1: let content quality speak for itself.",
        size=10, color=DGRAY)

    # Organic note
    y2 = y + i(0.85)
    rect(s, lx, y2, lw, i(0.44), LGRAY)
    txt(s, lx + i(0.15), y2 + i(0.07), lw - i(0.2), i(0.3),
        "95% Organic Distribution — no paid boosts, no promotions.",
        size=10, color=DARK)

    # Right divider
    rect(s, i(5.1), i(1.7), i(0.01), i(5.3), XGRAY)

    # Right platform column
    rx = i(5.4)
    rw = i(7.5)
    ry = i(1.7)
    label(s, rx, ry - i(0.25), rw, i(0.22), "Platform Execution")

    platforms = [
        ("Instagram  @yujiaaa_xie", IGRED,
         "Photo carousels with music pairing. Fashion + lifestyle in NYC cultural spaces (The Met). Captions in English with cultural references.",
         ["Lifestyle", "Fashion", "Carousel"]),
        ("Douyin  @Xyjjj", DARK,
         "Short-form videos with poetic, emotionally resonant Chinese captions. Aesthetic moments — art museums, quiet nights, candlelit scenes.",
         ["Emotional", "Poetic", "Short-form"]),
        ("Xiaohongshu  谢雨佳", XHSRED,
         "Lifestyle notes with location context. Same visual assets cross-distributed from Instagram — NYC life as the hook for Chinese audiences.",
         ["NYC Lifestyle", "Cross-distributed"]),
    ]

    for pname, pcolor, pdesc, ptags in platforms:
        rect(s, rx, ry, i(0.06), i(0.9), pcolor)
        rect(s, rx + i(0.06), ry, rw - i(0.06), i(0.9), BGRAY)
        txt(s, rx + i(0.22), ry + i(0.07), rw - i(0.3), i(0.26),
            pname, size=13, bold=True, color=DARK)
        txt(s, rx + i(0.22), ry + i(0.35), rw - i(0.3), i(0.45),
            pdesc, size=10, color=DGRAY)
        # Tags
        tx = rx + i(0.22)
        ty2 = ry + i(0.8)
        for tag_t in ptags:
            tw2 = i(len(tag_t) * 0.085 + 0.4)
            rect(s, tx, ty2, tw2, i(0.22), CREAM)
            txt(s, tx + i(0.07), ty2 + i(0.02), tw2, i(0.2),
                tag_t, size=9, bold=True, color=RGBColor(0xA0, 0x78, 0x40))
            tx += tw2 + i(0.1)
        ry += i(1.1)

    slide_num_tag(s, 4)
    return s


def slide5_baseline(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    section_header(s, "05  /  Before Activation", "Where We Started",
                   "Day 0 Snapshot — baseline metrics before any content was posted")

    platforms_data = [
        {
            "name": "Instagram", "handle": "@yujiaaa_xie",
            "img_idx": 2,
            "stats": [("2,387", "Followers"), ("22", "Interactions / 30 days")],
            "chip": "Established — Low Activity", "chip_c": LGREEN, "chip_t": GREEN,
        },
        {
            "name": "Xiaohongshu", "handle": "谢雨佳 · 423553458",
            "img_idx": 3,
            "stats": [("530", "Fans"), ("1,892", "Total Likes & Saves")],
            "chip": "Dormant — No Recent Posts", "chip_c": LAMBER, "chip_t": AMBER,
        },
        {
            "name": "Douyin", "handle": "@Xyjjj",
            "img_idx": 4,
            "stats": [("13", "Fans"), ("100", "Total Likes")],
            "chip": "Just Starting", "chip_c": RGBColor(0xE8, 0xEE, 0xFF), "chip_t": INDIG,
        },
    ]

    col_w = i(3.8)
    gap = i(0.4)
    total_w = col_w * 3 + gap * 2
    sx = (SW - total_w) / 2
    sy = i(1.65)

    for ci, pd in enumerate(platforms_data):
        cx = sx + ci * (col_w + gap)
        rect(s, cx, sy, col_w, i(5.2), BGRAY)

        # Screenshot
        ih = i(2.4)
        iw = col_w - i(0.3)
        ix = cx + i(0.15)
        iy = sy + i(0.15)
        img_buf = imgs[pd["img_idx"]]
        img_buf.seek(0)
        pil = Image.open(img_buf)
        aspect = pil.width / pil.height
        img_buf.seek(0)
        if aspect < 1:
            actual_w = ih * aspect
            actual_x = ix + (iw - actual_w) / 2
            pic(s, img_buf, actual_x, iy, actual_w, ih)
        else:
            actual_h = iw / aspect
            pic(s, img_buf, ix, iy, iw, actual_h)

        # Platform name
        wy = sy + i(2.65)
        txt(s, cx + i(0.2), wy, col_w - i(0.4), i(0.3),
            pd["name"], size=14, bold=True, color=DARK, font=HEAD_FONT)
        txt(s, cx + i(0.2), wy + i(0.32), col_w - i(0.4), i(0.2),
            pd["handle"], size=10, color=MGRAY)

        # Stats
        sy2 = wy + i(0.62)
        for val, lbl_t in pd["stats"]:
            txt(s, cx + i(0.2), sy2, col_w - i(0.4), i(0.4),
                val, size=26, bold=True, color=DARK, font=HEAD_FONT)
            label(s, cx + i(0.2), sy2 + i(0.42), col_w - i(0.4), i(0.2), lbl_t)
            sy2 += i(0.85)

        # Status chip
        rect(s, cx + i(0.2), sy + i(4.75), col_w - i(0.4), i(0.26), pd["chip_c"])
        txt(s, cx + i(0.3), sy + i(4.78), col_w - i(0.4), i(0.22),
            pd["chip"], size=10, bold=True, color=pd["chip_t"])

    slide_num_tag(s, 5)
    return s


def slide6_growth(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    section_header(s, "06  /  After Activation", "Growth Across All Three Platforms")

    # img[5]=IG after(800×558), img[6]=XHS after(368×800), img[7]=Douyin after(368×800)
    # img[8]=IG stats(800×493), img[9]=XHS stats(567×800)

    growth_data = [
        {
            "name": "Instagram", "color": IGRED, "img_idx": 5,
            "metrics": [
                ("Interactions / 30d", "22", "207", "+185"),
                ("Views / 30d", "14.4K", "21.7K", "+51%"),
                ("Followers", "2,387", "2,378", "-9 cleaned"),
            ],
        },
        {
            "name": "Xiaohongshu", "color": XHSRED, "img_idx": 6,
            "metrics": [
                ("Fans", "530", "556", "+26"),
                ("Likes & Saves", "1,892", "1,964", "+72"),
                ("Top Post Views", "—", "1,130", "new"),
            ],
        },
        {
            "name": "Douyin", "color": DARK, "img_idx": 7,
            "metrics": [
                ("Fans", "13", "267", "+254"),
                ("Total Likes", "100", "1,643", "+1,543"),
                ("Posts Published", "2", "10", "+8"),
            ],
        },
    ]

    col_w = i(3.8)
    gap = i(0.4)
    total_w = col_w * 3 + gap * 2
    sx = (SW - total_w) / 2
    sy = i(1.65)

    for ci, gd in enumerate(growth_data):
        cx = sx + ci * (col_w + gap)
        rect(s, cx, sy, col_w, i(4.4), BGRAY)
        rect(s, cx, sy, col_w, i(0.06), gd["color"])

        # Screenshot (small)
        ih = i(1.55)
        iw = col_w - i(0.3)
        ix = cx + i(0.15)
        iy = sy + i(0.15)
        img_buf = imgs[gd["img_idx"]]
        img_buf.seek(0)
        pil = Image.open(img_buf)
        aspect = pil.width / pil.height
        img_buf.seek(0)
        if aspect < 1:
            actual_w = ih * aspect
            actual_x = ix + (iw - actual_w) / 2
            pic(s, img_buf, actual_x, iy, actual_w, ih)
        else:
            actual_h = iw / aspect
            pic(s, img_buf, ix, iy, iw, actual_h)

        txt(s, cx + i(0.2), sy + i(1.8), col_w - i(0.4), i(0.3),
            gd["name"], size=13, bold=True, color=DARK, font=HEAD_FONT)

        ry = sy + i(2.15)
        for metric, before, after, delta in gd["metrics"]:
            label(s, cx + i(0.2), ry, col_w - i(0.4), i(0.18), metric)
            # before → after
            bx = cx + i(0.2)
            txt(s, bx, ry + i(0.2), i(0.8), i(0.32),
                before, size=13, color=RGBColor(0xBB, 0xBB, 0xBB), bold=True)
            txt(s, bx + i(0.82), ry + i(0.22), i(0.2), i(0.28),
                "→", size=11, color=GOLD)
            txt(s, bx + i(1.05), ry + i(0.2), i(1.0), i(0.32),
                after, size=16, bold=True, color=DARK, font=HEAD_FONT)
            # delta chip
            dc = LGREEN if delta.startswith('+') else LAMBER if delta.startswith('-') else LGRAY
            dtc = GREEN if delta.startswith('+') else AMBER if delta.startswith('-') else MGRAY
            rect(s, bx + i(2.15), ry + i(0.24), i(0.9), i(0.24), dc)
            txt(s, bx + i(2.2), ry + i(0.26), i(0.9), i(0.2),
                delta, size=9, bold=True, color=dtc)
            ry += i(0.68)

    # KPI banner
    by = sy + i(4.6)
    rect(s, i(0.5), by, SW - i(1.0), i(0.7), DARK)
    txt(s, i(0.8), by + i(0.12), i(4), i(0.25),
        "TOTAL VALID KPIs  —  TARGET: 500", size=9, bold=True,
        color=RGBColor(0x55, 0x55, 0x55))
    txt(s, i(4.5), by + i(0.08), i(2.5), i(0.5),
        "2,277", size=36, bold=True, color=GOLD, font=HEAD_FONT)
    txt(s, i(7.1), by + i(0.15), i(5), i(0.35),
        "4.6× over target  ·  in under 1 week", size=13, color=WHITE)

    slide_num_tag(s, 6)
    return s


def slide7_worked(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    section_header(s, "07  /  What Worked",
                   "Two Hero Posts. Two Platforms. One Strategy.")

    card_y = i(1.65)
    card_h = i(4.85)
    gap = i(0.3)
    card_w = (SW - i(1.4) - gap) / 2
    c1x = i(0.7)
    c2x = c1x + card_w + gap

    # Card 1 — Douyin
    rect(s, c1x, card_y, card_w, card_h, DARK)

    txt(s, c1x + i(0.25), card_y + i(0.2), card_w - i(0.5), i(0.9),
        "28K", size=56, bold=True, color=GOLD, font=HEAD_FONT)
    txt(s, c1x + i(0.25), card_y + i(1.0), card_w - i(0.5), i(0.22),
        "PLAYS · HERO POST", size=9, bold=True, color=MGRAY)

    # Douyin screenshot – img[10] (322×700)
    ph = i(2.3)
    pw = ph * (322 / 700)
    sx2 = c1x + (card_w - pw) / 2
    pic(s, imgs[10], sx2, card_y + i(1.3), pw, ph)

    txt(s, c1x + i(0.25), card_y + i(3.7), card_w - i(0.5), i(0.22),
        "DOUYIN  @Xyjjj", size=9, bold=True, color=GOLD)
    txt(s, c1x + i(0.25), card_y + i(3.95), card_w - i(0.5), i(0.4),
        "我擅长感受 却拙于表达", size=14, bold=True, color=WHITE)
    txt(s, c1x + i(0.25), card_y + i(4.38), card_w - i(0.5), i(0.22),
        "I Feel Deeply, Struggle to Express", size=10, italic=True,
        color=RGBColor(0x88, 0x88, 0x88))

    # Metrics chips
    metrics1 = ["900 Likes", "23 Comments", "94 Saves", "+155 Followers"]
    mx = c1x + i(0.25)
    my = card_y + i(4.65)
    for m in metrics1:
        mw = i(len(m) * 0.085 + 0.3)
        rect(s, mx, my, mw, i(0.22), RGBColor(0x2A, 0x2A, 0x2A))
        txt(s, mx + i(0.07), my + i(0.02), mw, i(0.2),
            m, size=9, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC))
        mx += mw + i(0.08)

    # Card 2 — Instagram (Met Museum)
    rect(s, c2x, card_y, card_w, card_h, BGRAY)

    # Met Museum photo – img[11] (1319×989)
    iw2 = card_w
    ih2 = iw2 * (989 / 1319)
    pic(s, imgs[11], c2x, card_y, iw2, ih2)

    iy2_base = card_y + ih2 + i(0.15)
    txt(s, c2x + i(0.25), iy2_base, card_w - i(0.5), i(0.22),
        "INSTAGRAM  @yujiaaa_xie", size=9, bold=True, color=IGRED)
    txt(s, c2x + i(0.25), iy2_base + i(0.25), card_w - i(0.5), i(0.3),
        "The Met Museum — NYC", size=14, bold=True, color=DARK, font=HEAD_FONT)
    txt(s, c2x + i(0.25), iy2_base + i(0.6), card_w - i(0.5), i(0.22),
        "Cultural space · editorial fashion · golden-hour light", size=10, color=DGRAY)

    metrics2 = ["110 Likes", "12 Comments", "2.7K Views", "141 Profile Visits"]
    mx2 = c2x + i(0.25)
    my2 = iy2_base + i(0.9)
    for m in metrics2:
        mw = i(len(m) * 0.085 + 0.3)
        rect(s, mx2, my2, mw, i(0.22), CREAM)
        txt(s, mx2 + i(0.07), my2 + i(0.02), mw, i(0.2),
            m, size=9, bold=True, color=RGBColor(0xA0, 0x78, 0x40))
        mx2 += mw + i(0.08)

    # Key insight bar
    ib_y = card_y + card_h + i(0.12)
    rect(s, i(0.7), ib_y, SW - i(1.4), i(0.6), CREAM)
    txt(s, i(0.95), ib_y + i(0.1), SW - i(1.9), i(0.42),
        "Key insight:  The same Met Museum shoot produced the #1 post on both Instagram and Xiaohongshu simultaneously — "
        "one creative execution, two platform wins. Quality of location + subject > quantity of posts.",
        size=11, color=DARK)

    slide_num_tag(s, 7)
    return s


def slide8_didnt(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    section_header(s, "08  /  What Didn't Work", "Three Gaps. Three Lessons.")

    issues = [
        {
            "icon": "📉", "title": "Douyin — Inconsistent Reach",
            "top_color": RGBColor(0xE5, 0x39, 0x35),
            "desc": ("Post 1 hit 24K plays. Posts 2 and 3 dropped to 3.4K and 1.5K respectively. "
                     "The algorithm rewarded the emotional hook the first time — but follow-up content "
                     "couldn't sustain that signal."),
            "stat": "Posts 2+3 avg: 2,421 plays vs 24K hero",
        },
        {
            "icon": "🔍", "title": "Xiaohongshu — Discovery Gap",
            "top_color": RGBColor(0xF5, 0x9E, 0x0B),
            "desc": ("The Met post got 1,130 views and a 6.2% interaction rate — but Posts 2 and 3 "
                     "collapsed to 144 and 201 views. Without hashtag strategy, Xiaohongshu's "
                     "discovery layer didn't amplify past the first post."),
            "stat": "Posts 2+3 avg: 172 views vs 1,130 hero",
        },
        {
            "icon": "🏷️", "title": "No-Hashtag Trade-off",
            "top_color": INDIG,
            "desc": ("The deliberate no-hashtag bet worked on Douyin (algorithm-driven) but limited "
                     "reach on Xiaohongshu, where topic tags are the primary discovery mechanism. "
                     "Pure organic content quality was not enough to beat the platform's tag-based feed."),
            "stat": "Phase 2 hypothesis: add 3–5 topic tags",
        },
    ]

    col_w = i(3.8)
    gap = i(0.4)
    total_w = col_w * 3 + gap * 2
    sx = (SW - total_w) / 2
    sy = i(1.65)
    card_h = i(4.1)

    for ci, issue in enumerate(issues):
        cx = sx + ci * (col_w + gap)
        rect(s, cx, sy, col_w, card_h, BGRAY)
        rect(s, cx, sy, col_w, i(0.07), issue["top_color"])

        txt(s, cx + i(0.22), sy + i(0.22), i(0.5), i(0.4), issue["icon"], size=22)
        txt(s, cx + i(0.22), sy + i(0.65), col_w - i(0.44), i(0.4),
            issue["title"], size=14, bold=True, color=DARK, font=HEAD_FONT)
        txt(s, cx + i(0.22), sy + i(1.1), col_w - i(0.44), i(2.0),
            issue["desc"], size=11, color=DGRAY)

        rect(s, cx + i(0.22), sy + card_h - i(0.52), col_w - i(0.44), i(0.3), XGRAY)
        txt(s, cx + i(0.32), sy + card_h - i(0.48), col_w - i(0.5), i(0.26),
            issue["stat"], size=10, bold=True, color=MGRAY)

    # Phase 2 adjustment bar
    ab_y = sy + card_h + i(0.22)
    rect(s, i(0.7), ab_y, SW - i(1.4), i(0.6), DARK)
    txt(s, i(0.95), ab_y + i(0.08), i(1.8), i(0.3),
        "ADJUSTMENT →", size=9, bold=True, color=GOLD)
    txt(s, i(2.95), ab_y + i(0.08), SW - i(3.7), i(0.4),
        "Phase 2: tighter emotional brief for Douyin consistency  ·  topic hashtags on Xiaohongshu  ·  replicate Met-style cultural shoots on Instagram",
        size=11, color=RGBColor(0xAA, 0xAA, 0xAA))

    slide_num_tag(s, 8)
    return s


def slide9_kpi(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, DARK2)
    accent_bar(s)
    project_tag(s, dark=True)

    # Top label
    txt(s, 0, i(1.0), SW, i(0.3),
        "VALID KPIs  ·  LIKES + COMMENTS + SHARES + SAVES + NEW FOLLOWERS",
        size=10, bold=True, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

    # Big number
    txt(s, 0, i(1.5), SW, i(2.5),
        "2,277", size=120, bold=True, color=GOLD, font=HEAD_FONT, align=PP_ALIGN.CENTER)

    txt(s, 0, i(3.9), SW, i(0.4),
        "4.6×  over the 500-person target  ·  in under 1 week",
        size=15, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

    # KPI breakdown
    kpis = [
        ("1,849", "Likes & Saves"),
        ("64", "Comments"),
        ("64", "Shares"),
        ("300", "New Followers"),
    ]
    kw = SW / len(kpis)
    ky = i(4.55)
    rect(s, 0, ky - i(0.02), SW, i(0.01), RGBColor(0x1E, 0x1E, 0x1E))

    for idx, (kval, klbl) in enumerate(kpis):
        kx = idx * kw
        if idx > 0:
            rect(s, kx, ky, i(0.01), i(1.4), RGBColor(0x1E, 0x1E, 0x1E))
        txt(s, kx, ky + i(0.15), kw, i(0.7),
            kval, size=32, bold=True, color=WHITE, font=HEAD_FONT, align=PP_ALIGN.CENTER)
        txt(s, kx, ky + i(0.85), kw, i(0.3),
            klbl.upper(), size=9, bold=True, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

    # Platform breakdown
    pf_data = [
        ("1,974", "Douyin", "86% of total"),
        ("201", "Instagram", "9% of total"),
        ("102", "Xiaohongshu", "5% of total"),
    ]
    pfy = i(6.0)
    pw2 = SW / 3
    rect(s, 0, pfy - i(0.01), SW, i(0.01), RGBColor(0x1E, 0x1E, 0x1E))
    for idx, (pval, pname, ppct) in enumerate(pf_data):
        px = idx * pw2
        if idx > 0:
            rect(s, px, pfy, i(0.01), i(1.3), RGBColor(0x1E, 0x1E, 0x1E))
        txt(s, px, pfy + i(0.1), pw2, i(0.45),
            pval, size=22, bold=True, color=WHITE, font=HEAD_FONT, align=PP_ALIGN.CENTER)
        txt(s, px, pfy + i(0.55), pw2, i(0.25),
            pname, size=11, bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
        txt(s, px, pfy + i(0.82), pw2, i(0.22),
            ppct, size=10, color=GOLD, align=PP_ALIGN.CENTER)

    slide_num_tag(s, 9)
    return s


def slide10_learnings(prs, imgs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    set_bg(s, WHITE)
    accent_bar(s)
    project_tag(s)

    section_header(s, "10  /  Learnings",
                   "Four Insights That Will Shape Phase 2")

    cards = [
        {
            "num": "01  —  KEY LEARNING", "dark": True,
            "title": "Emotion > Aesthetics on Chinese Platforms",
            "text": ('Douyin\'s viral post succeeded not because of beautiful visuals — but because the caption '
                     '"我擅长感受 却拙于表达" created emotional identification. Chinese platform audiences reward '
                     'resonance over production quality. The algorithm follows the feeling.'),
        },
        {
            "num": "02", "dark": False,
            "title": "One Shoot = Three Audiences",
            "text": ("The Met Museum shoot became the top post on both Instagram and Xiaohongshu simultaneously. "
                     "Cross-platform repurposing isn't just efficient — it's strategic. The right shoot, edited "
                     "for each platform's tone, compounds its own return."),
        },
        {
            "num": "03", "dark": False,
            "title": "Content Quality Beats Tactics",
            "text": ("No paid promotion. No hashtags. No collaborations. 2,277 KPIs came purely from posting the "
                     "right content in the right emotional frame. Phase 1 confirmed: when content is strong "
                     "enough, the platform works for you."),
        },
        {
            "num": "04", "dark": False,
            "title": "Platform Culture Shapes What Lands",
            "text": ("The same no-hashtag strategy worked on Douyin (algorithm-first) but failed on Xiaohongshu "
                     "(tag-discovery-first). Understanding each platform's native discovery logic — not just its "
                     "aesthetic — is the difference between 1,130 views and 144."),
        },
    ]

    cw = (SW - i(1.4) - i(0.3)) / 2
    ch = (SH - i(2.0) - i(0.3)) / 2
    positions = [
        (i(0.7), i(1.65)),
        (i(0.7) + cw + i(0.3), i(1.65)),
        (i(0.7), i(1.65) + ch + i(0.3)),
        (i(0.7) + cw + i(0.3), i(1.65) + ch + i(0.3)),
    ]

    for ci, (card, (cx, cy)) in enumerate(zip(cards, positions)):
        bg_c = DARK if card["dark"] else BGRAY
        top_c = GOLD if card["dark"] else XGRAY
        num_c = RGBColor(0x55, 0x55, 0x55) if card["dark"] else MGRAY
        title_c = GOLD if card["dark"] else DARK
        text_c = RGBColor(0xAA, 0xAA, 0xAA) if card["dark"] else DGRAY

        rect(s, cx, cy, cw, ch, bg_c)
        rect(s, cx, cy, cw, i(0.06), top_c)

        txt(s, cx + i(0.25), cy + i(0.2), cw - i(0.5), i(0.22),
            card["num"], size=9, bold=True, color=num_c)
        txt(s, cx + i(0.25), cy + i(0.46), cw - i(0.5), i(0.5),
            card["title"], size=14, bold=True, color=title_c, font=HEAD_FONT)
        txt(s, cx + i(0.25), cy + i(1.0), cw - i(0.5), ch - i(1.1),
            card["text"], size=11, color=text_c)

    slide_num_tag(s, 10)
    return s


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    html_path = '/Users/wangbo/Desktop/NYU/Third semaster/Social Media & Brand/@Yujia Activation/@Yujia presentation final ver.2.html'
    out_path  = '/Users/wangbo/Desktop/NYU/Third semaster/Social Media & Brand/@Yujia Activation/East_Meets_Feed.pptx'

    print("Extracting images...")
    imgs = extract_images(html_path)
    print(f"  {len(imgs)} images extracted")

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides...")
    slide1_title(prs, imgs)
    print("  Slide 1 — Title")
    slide2_subject(prs, imgs)
    print("  Slide 2 — The Subject & Why")
    slide3_strategy(prs, imgs)
    print("  Slide 3 — Platform Strategy")
    slide4_plan(prs, imgs)
    print("  Slide 4 — Activation Plan")
    slide5_baseline(prs, imgs)
    print("  Slide 5 — Before Activation")
    slide6_growth(prs, imgs)
    print("  Slide 6 — After Activation")
    slide7_worked(prs, imgs)
    print("  Slide 7 — What Worked")
    slide8_didnt(prs, imgs)
    print("  Slide 8 — What Didn't Work")
    slide9_kpi(prs, imgs)
    print("  Slide 9 — KPI Proof")
    slide10_learnings(prs, imgs)
    print("  Slide 10 — Learnings")

    prs.save(out_path)
    print(f"\nSaved → {out_path}")


if __name__ == "__main__":
    main()
